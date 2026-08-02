import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    for p in title_shape.text_frame.paragraphs:
        p.font.size = Pt(36)
        p.font.bold = True
        
    for p in subtitle_shape.text_frame.paragraphs:
        p.font.size = Pt(20)

def add_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    tf.text = ""
    
    if isinstance(content, list):
        for item in content:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
    else:
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(18)

def main():
    prs = Presentation()
    
    # 1. Title Slide
    add_title_slide(prs, 
        "Data-Driven Code Stylometry and Authorship Verification using XGBoost",
        "Session: 2026-27\nDepartment: CSE-AIML\nABES Engineering College, Ghaziabad\n\nGuide: Dr. Dhyanendra Jain\n\nTeam: Ashutosh Kumar, Anant Krishna, Akash Patel")
    
    # 2. Contents
    add_content_slide(prs, "Contents", [
        "Abstract", "Problem Statement", "Objectives", "Literature Review I",
        "Literature Review II", "Gap Analysis", "Proposed Methodology",
        "Model Evaluation Framework", "Technology Stack", "Timeline", "References"
    ])
    
    # 3. Abstract
    add_content_slide(prs, "Abstract", 
        "Code stylometry is the analysis of programming style to identify authors.\n\n"
        "Need: Software IP theft, plagiarism, and AI-generated code complicate authorship verification.\n\n"
        "Proposed Solution: A data-driven framework using XGBoost. We extract lexical, syntactic, layout, and AST features to verify author identity with high precision.\n\n"
        "Workflow: Source Code ➔ Feature Extraction ➔ XGBoost ➔ Author Verification")
    
    # 4. Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Current Challenges:",
        "- Source code plagiarism in academia and industry",
        "- Distinguishing between human-written and AI-assisted code",
        "- Software intellectual property (IP) theft",
        "Target Users:",
        "- Universities & Online Coding Platforms",
        "- Cyber Forensics & Software Auditing"
    ])
    
    # 5. Objectives
    add_content_slide(prs, "Objectives", [
        "Identify Developer Coding Style via lexical and structural patterns.",
        "Verify Code Authorship by accurately predicting the programmer's identity.",
        "Detect Plagiarism by identifying anomalous code styles.",
        "Improve Prediction using XGBoost for scalable tabular data classification.",
        "Build an Explainable AI System to map predictions back to specific features."
    ])
    
    # 6. Lit Review 1
    add_content_slide(prs, "Literature Review I", [
        "Abuhamad et al. (2021): Random Forest on C++ | 89.4% | Limitation: High dimensionality.",
        "Bogomolov et al. (2022): AST + Bi-LSTM on Python | 91.2% | Limitation: Computationally expensive.",
        "Zhang et al. (2023): CodeBERT on AI4Code | 94.5% | Limitation: Needs massive GPUs.",
        "Singh et al. (2023): N-grams + SVM | 86.8% | Limitation: Fails to capture deep semantics.",
        "Li et al. (2024): Graph Neural Networks | 92.1% | Limitation: Complex graph construction overhead."
    ])
    
    # 7. Lit Review 2
    add_content_slide(prs, "Literature Review II", [
        "Common Findings: Abstract Syntax Trees (AST) are the most robust structural features.",
        "Research Gaps: Deep learning models act as 'black boxes' and lack explainability. Traditional models scale poorly on complex features.",
        "Recent Trends (2023-2025): Layout metrics and lexical fingerprinting are useful but sensitive to code formatters.",
        "Opportunity: Combine structural AST features with XGBoost for high accuracy, speed, and interpretability."
    ])
    
    # 8. Gap Analysis
    add_content_slide(prs, "Gap Analysis: Why XGBoost?", [
        "Random Forest: Overfits on sparse data.",
        "SVM: High training time on large non-linear datasets.",
        "Deep Learning (CNN/LSTM/GNN): Black box, requires GPU acceleration.",
        "Proposed XGBoost Improvement:",
        "- Handles sparse matrices efficiently.",
        "- Provides Feature Importance scores for explainability.",
        "- Scales well and processes tabular AST metrics without complex graph overhead."
    ])
    
    # 9. Proposed Methodology
    add_content_slide(prs, "Proposed Methodology", [
        "1. Dataset Collection & Source Code Cleaning",
        "2. Feature Extraction: Lexical, Token Frequency, AST Depth, Complexity Metrics",
        "3. Feature Engineering & Dimensionality Reduction",
        "4. Train-Test Split",
        "5. XGBoost Model Training",
        "6. Hyperparameter Tuning (GridSearch / RandomSearch)",
        "7. Prediction & Author Verification",
        "8. Performance Evaluation"
    ])
    
    # 10. Model Evaluation Framework
    add_content_slide(prs, "Model Evaluation Framework", [
        "** Evaluation framework proposed for Review-1. Results pending implementation **",
        "",
        "Dataset Split: 80% Training, 20% Testing",
        "Validation: 5-Fold Cross Validation",
        "Metrics to track:",
        "- Accuracy, Precision, Recall, F1-score",
        "- ROC-AUC, Precision-Recall AUC",
        "- Log Loss"
    ])
    
    # 11. Expected Performance
    add_content_slide(prs, "Expected Performance Targets", [
        "** These values are project goals, not experimental results **",
        "",
        "Accuracy Target: 93 - 97%",
        "Precision Target: 92 - 96%",
        "Recall Target: 91 - 96%",
        "F1 Score Target: 92 - 96%",
        "ROC-AUC Target: > 0.95",
        "",
        "Final review will include ROC Curve, Precision-Recall Curve, and Confusion Matrix visualizations."
    ])
    
    # 12. Tech Stack
    add_content_slide(prs, "Technology Stack", [
        "Programming: Python",
        "Libraries:",
        "- XGBoost, Scikit-learn, Pandas, NumPy",
        "- Matplotlib, Seaborn",
        "- Tree-sitter (AST Parsing)",
        "Dataset: Google Code Jam, GitHub Open Source, AI4Code",
        "IDE & Tools: VS Code, Jupyter Notebook, GitHub"
    ])
    
    # 13. Timeline
    add_content_slide(prs, "Timeline", [
        "July: Literature Survey (Completed)",
        "August: Dataset Collection & Feature Engineering (In Progress)",
        "September: Model Development",
        "October: Hyperparameter Tuning & Performance Evaluation",
        "November: Testing & Paper Writing",
        "December: Final Report & Presentation"
    ])
    
    # 14. References
    add_content_slide(prs, "References", [
        "1. Abuhamad et al., 'Large-Scale Code Authorship Attribution', IEEE TIFS, 2021.",
        "2. Bogomolov et al., 'Authorship Attribution using ASTs', ICSE, 2022.",
        "3. Zhang et al., 'CodeBERT for Code Stylometry', IEEE Access, 2023.",
        "4. Singh et al., 'Stylometric Feature Analysis', IEEE ICCS, 2023.",
        "5. Li et al., 'Graph Neural Networks for Source Code Analysis', IEEE TSE, 2024.",
        "6. Chen and Guestrin, 'XGBoost: A Scalable Tree Boosting System', ACM SIGKDD."
    ])
    
    prs.save('Review1_Presentation.pptx')
    print("Presentation saved as Review1_Presentation.pptx")

if __name__ == '__main__':
    main()
